import copy
import io
import json
import logging
import os
import random
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from lxml import etree

log = logging.getLogger(__name__)

TEMPLATES_DIR = Path(__file__).resolve().parent.parent / 'data' / 'templates'
DEFAULT_PPTX = TEMPLATES_DIR / 'default.pptx'
PROTOTYPE_PATH = os.environ.get(
    'PPTX_TEMPLATE_PATH',
    str(Path(__file__).resolve().parent / 'template.pptx'),
)

TYPE_TO_PROTO_IDX = {
    "title": [0],
    "bullets": [1],
    "two_column": [2],
    "image_text": [3],
    "section": [4, 5, 6, 7, 8, 9],
    "table": [10],
    "thank_you": [11],
}
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _set_text_frame_text(tf, text, is_title=False, color=None):
    if tf is None:
        return
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(text)
    if is_title:
        p.font.size = Pt(36)
    else:
        p.font.size = Pt(18)
    if color:
        from pptx.dml.color import RGBColor
        p.font.color.rgb = RGBColor(*color)


def _add_bullets_to_textframe(tf, items):
    if tf is None or not items:
        return
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        p.font.size = Pt(16)


def _add_image(picture_ph, image_path, slide):
    if picture_ph is not None:
        try:
            picture_ph.insert_picture(image_path)
            return
        except Exception:
            pass
    try:
        left = Inches(1)
        top = Inches(2)
        height = Inches(4)
        slide.shapes.add_picture(image_path, left, top, height=height)
    except Exception as e:
        log.warning(f'Failed to insert image: {e}')


def _add_table(placeholder, table_data, slide):
    if not table_data or not table_data[0]:
        return
    rows = len(table_data)
    cols = len(table_data[0])
    rows = max(rows, 1)
    cols = max(cols, 1)

    if placeholder is not None:
        try:
            shape = placeholder.insert_table(rows, cols)
            table = shape.table
            for r_idx, row in enumerate(table_data):
                for c_idx, val in enumerate(row):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = str(val) if val else ''
            return
        except Exception:
            pass

    left = Inches(1)
    top = Inches(2)
    width = Inches(10)
    height = Inches(0.5) * rows
    try:
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        for r_idx, row in enumerate(table_data):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(val) if val else ''
    except Exception as e:
        log.warning(f'Failed to add table: {e}')


def _get_text_shapes(slide):
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_shapes.append(shape)
    return text_shapes


def _fill_title_slide(slide, spec):
    text_shapes = _get_text_shapes(slide)
    title_text = _clean_title(spec.get('title', ''))
    subtitle = spec.get('subtitle', '')
    white = (255, 255, 255)
    if text_shapes:
        _set_text_frame_text(text_shapes[0].text_frame, title_text, is_title=True, color=white)
    if subtitle and len(text_shapes) > 1:
        _set_text_frame_text(text_shapes[1].text_frame, subtitle, color=white)


def _filter_placeholder_items(items):
    return [i for i in items if not any(p in i for p in PROTOTYPE_PLACEHOLDER_TEXTS)]


def _clean_title(title):
    if any(p in title for p in PROTOTYPE_PLACEHOLDER_TEXTS):
        return ''
    return title


def _fill_bullets_slide(slide, spec):
    text_shapes = _get_text_shapes(slide)
    title_text = _clean_title(spec.get('title', ''))
    items = _filter_placeholder_items(spec.get('items', []))
    if text_shapes:
        _set_text_frame_text(text_shapes[0].text_frame, title_text, is_title=True)
    if items and len(text_shapes) > 1:
        _add_bullets_to_textframe(text_shapes[1].text_frame, items)


def _fill_two_column_slide(slide, spec):
    text_shapes = _get_text_shapes(slide)
    title_text = _clean_title(spec.get('title', ''))
    left_text = spec.get('left', '')
    right_text = spec.get('right', '')
    if text_shapes:
        _set_text_frame_text(text_shapes[0].text_frame, title_text, is_title=True)
    if len(text_shapes) > 1:
        _set_text_frame_text(text_shapes[1].text_frame, left_text)
    if len(text_shapes) > 2:
        _set_text_frame_text(text_shapes[2].text_frame, right_text)


def _fill_image_text_slide(slide, spec, image_paths):
    text_shapes = _get_text_shapes(slide)
    title_text = _clean_title(spec.get('title', ''))
    image_fid = spec.get('image_file_id', '')
    caption = spec.get('image_caption', '')
    if text_shapes:
        _set_text_frame_text(text_shapes[0].text_frame, title_text, is_title=True)

    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    picture_ph = None
    for ph in phs.values():
        if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            picture_ph = ph
            break

    if image_fid and image_fid in image_paths:
        _add_image(picture_ph, image_paths[image_fid], slide)
    elif image_fid and image_fid in image_paths:
        _add_image(None, image_paths[image_fid], slide)

    if caption and len(text_shapes) > 1:
        _set_text_frame_text(text_shapes[1].text_frame, caption)


def _fill_section_slide(slide, spec):
    text_shapes = _get_text_shapes(slide)
    title_text = _clean_title(spec.get('title', ''))
    subtitle = spec.get('subtitle', '')
    white = (255, 255, 255)
    if text_shapes:
        _set_text_frame_text(text_shapes[0].text_frame, title_text, is_title=True, color=white)
    if subtitle and len(text_shapes) > 1:
        _set_text_frame_text(text_shapes[1].text_frame, subtitle, color=white)


PROTOTYPE_PLACEHOLDER_TEXTS = {'НАЗВАНИЕ', 'ПРЕЗЕНТАЦИИ', 'СЛАЙД', 'ПЕРЕБИВОЧНЫЙ', 'СПАСИБО'}


def _any_placeholder_cell(table_data):
    for row in table_data:
        for cell in row:
            if cell:
                for placeholder in PROTOTYPE_PLACEHOLDER_TEXTS:
                    if placeholder in cell:
                        return True
    return False


def _fill_table_slide(slide, spec):
    text_shapes = _get_text_shapes(slide)
    title_text = _clean_title(spec.get('title', ''))
    if text_shapes:
        _set_text_frame_text(text_shapes[0].text_frame, title_text, is_title=True)

    table_data = spec.get('table_data', None)
    if table_data is None:
        header = spec.get('header', [])
        rows = spec.get('rows', [])
        if header:
            table_data = [header] + rows
        else:
            table_data = rows or []

    if not table_data or not table_data[0]:
        return

    if table_data and not table_data[0][0]:
        first_col_vals = [row[0] for row in table_data if row and row[0]]
        if len(first_col_vals) > 1:
            items = [v for v in first_col_vals
                     if not any(p in v for p in PROTOTYPE_PLACEHOLDER_TEXTS)]
            if items and len(text_shapes) > 1:
                _add_bullets_to_textframe(text_shapes[1].text_frame, items)
            return

    if table_data[0] and len(table_data[0]) == 1:
        items = [row[0] for row in table_data if row and row[0]
                 and not any(p in row[0] for p in PROTOTYPE_PLACEHOLDER_TEXTS)]
        if not items:
            return
        if items and len(text_shapes) > 1:
            _add_bullets_to_textframe(text_shapes[1].text_frame, items)
        return

    if _any_placeholder_cell(table_data):
        items = []
        for row in table_data:
            for cell in row:
                if cell and not any(p in cell for p in PROTOTYPE_PLACEHOLDER_TEXTS):
                    items.append(cell)
        if items and len(text_shapes) > 1:
            _add_bullets_to_textframe(text_shapes[1].text_frame, items)
        return

    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    content_ph = phs.get(0) or phs.get(1) if phs else None
    _add_table(content_ph, table_data, slide)


def _fill_thank_you_slide(slide, spec):
    text_shapes = _get_text_shapes(slide)
    title_text = _clean_title(spec.get('title', '')) or 'СПАСИБО!'
    if text_shapes:
        _set_text_frame_text(text_shapes[0].text_frame, title_text, is_title=True, color=(255, 255, 255))


FILL_DISPATCH = {
    "title": _fill_title_slide,
    "bullets": _fill_bullets_slide,
    "two_column": _fill_two_column_slide,
    "image_text": _fill_image_text_slide,
    "section": _fill_section_slide,
    "table": _fill_table_slide,
    "thank_you": _fill_thank_you_slide,
}


def _duplicate_slide(prs, slide_index):
    source = prs.slides[slide_index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)

    rId_map = {}
    for rel_key, rel in source.part.rels.items():
        reltype = rel.reltype
        if 'slideLayout' in reltype or 'slideMaster' in reltype:
            continue
        if not rel.is_external:
            try:
                new_rId = new_slide.part.rels.get_or_add(reltype, rel.target_part)
                rId_map[rel_key] = new_rId
            except Exception:
                pass

    spTree = new_slide.shapes._spTree
    for child in list(spTree):
        tag = child.tag
        if tag.endswith('}sp') or tag.endswith('}pic') or \
           tag.endswith('}grpSp') or tag.endswith('}graphicFrame'):
            spTree.remove(child)

    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    source_cSld = source.element.find(f'{{{ns}}}cSld')
    new_cSld = new_slide.element.find(f'{{{ns}}}cSld')

    for shape in source.shapes:
        cloned = copy.deepcopy(shape._element)
        for elem in cloned.iter():
            for attr in ('embed', 'link'):
                val = elem.get(f'{{{r_ns}}}{attr}')
                if val and val in rId_map:
                    elem.set(f'{{{r_ns}}}{attr}', rId_map[val])
        spTree.append(cloned)

    if source_cSld is not None and new_cSld is not None:
        source_bg = source_cSld.find(f'{{{ns}}}bg')
        existing_bg = new_cSld.find(f'{{{ns}}}bg')
        if existing_bg is not None:
            new_cSld.remove(existing_bg)
        if source_bg is not None:
            new_bg = copy.deepcopy(source_bg)
            for elem in new_bg.iter():
                for attr in ('embed', 'link'):
                    val = elem.get(f'{{{r_ns}}}{attr}')
                    if val and val in rId_map:
                        elem.set(f'{{{r_ns}}}{attr}', rId_map[val])
            new_cSld.insert(0, new_bg)

    return len(list(prs.slides)) - 1


def _remove_slide(prs, slide_index):
    pres_elem = prs.part.presentation.element
    sldIdLst = pres_elem.sldIdLst
    entry = sldIdLst[slide_index]
    rId = entry.get(f'{{{NS_R}}}id')
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
    del sldIdLst[slide_index]


def _reorder_slides(prs, new_order):
    pres_elem = prs.part.presentation.element
    sldIdLst = pres_elem.sldIdLst
    els = list(sldIdLst)
    reordered = [els[i] for i in new_order]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in reordered:
        sldIdLst.append(el)


async def generate_pptx(slides, image_paths=None, use_image_placeholders=True):
    image_paths = image_paths or {}

    if not os.path.isfile(PROTOTYPE_PATH):
        raise FileNotFoundError(f'Template not found: {PROTOTYPE_PATH}')

    prs = Presentation(PROTOTYPE_PATH)

    type_count = {}
    for s in slides:
        t = s.get('type', 'bullets')
        type_count[t] = type_count.get(t, 0) + 1

    type_available = {}
    for t, count in type_count.items():
        if count == 0:
            continue
        candidates = TYPE_TO_PROTO_IDX.get(t)
        if not candidates:
            continue
        indices = []
        if len(candidates) > 1:
            pool = list(candidates)
            random.shuffle(pool)
            for i in range(count):
                if not pool:
                    pool = list(candidates)
                    random.shuffle(pool)
                proto = pool.pop(0)
                if i == 0:
                    indices.append(proto)
                else:
                    new_idx = _duplicate_slide(prs, proto)
                    indices.append(new_idx)
        else:
            proto = candidates[0]
            indices = [proto]
            for _ in range(count - 1):
                indices.append(_duplicate_slide(prs, proto))
        type_available[t] = indices

    type_pools = {t: list(indices) for t, indices in type_available.items()}
    output_proto_indices = []
    for s in slides:
        t = s.get('type', 'bullets')
        pool = type_pools.get(t)
        if pool:
            idx = pool.pop(0)
            output_proto_indices.append(idx)

    used_indices = set(output_proto_indices)
    all_count = len(list(prs.slides))
    for idx in range(all_count - 1, -1, -1):
        if idx not in used_indices:
            try:
                _remove_slide(prs, idx)
            except Exception as e:
                log.warning(f'Failed to remove slide {idx}: {e}')

    surviving = sorted(used_indices)
    old_to_new = {old: new for new, old in enumerate(surviving)}
    desired_new_order = [old_to_new[idx] for idx in output_proto_indices]
    try:
        _reorder_slides(prs, desired_new_order)
    except Exception as e:
        log.warning(f'Failed to reorder slides: {e}')

    for slide_idx, slide_spec in zip(range(len(prs.slides)), slides):
        try:
            _fill_slide_prototype(prs.slides[slide_idx], slide_spec, image_paths)
        except Exception as e:
            log.warning(f'Failed to fill slide {slide_idx}: {e}')

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _fill_slide_prototype(slide, spec, image_paths=None):
    image_paths = image_paths or {}
    slide_type = spec.get('type', 'bullets')
    filler = FILL_DISPATCH.get(slide_type)
    if filler:
        if slide_type == 'image_text':
            filler(slide, spec, image_paths)
        else:
            filler(slide, spec)
